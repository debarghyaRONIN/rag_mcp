import os
import uuid
import json
import logging
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
from datetime import datetime
from enum import Enum

# FastAPI imports
from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.responses import JSONResponse
from pydantic import BaseModel

# Document processing imports
import PyPDF2
from docx import Document
from pptx import Presentation
import csv
import pandas as pd

# LangChain imports
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document as LangChainDocument

# Vector store and embeddings
import chromadb
from sentence_transformers import SentenceTransformer
import numpy as np

# HTTP client for Ollama
import requests

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize FastAPI app
app = FastAPI(title="RAG Chatbot with Agentic Architecture & Chain of Thought", version="1.0.0")

# Global configurations
CHROMA_DB_PATH = "./chroma_db"
OLLAMA_BASE_URL = "http://localhost:11434"
EMBEDDING_MODEL = "all-MiniLM-L6-v2"
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
TOP_K_RETRIEVAL = 5

# Chain of Thought reasoning types
class ReasoningType(Enum):
    SIMPLE = "simple"
    COMPLEX = "complex"
    ANALYTICAL = "analytical"
    COMPARATIVE = "comparative"
    CREATIVE = "creative"

# Chain of Thought step
class CoTStep(BaseModel):
    step_number: int
    reasoning: str
    evidence: Optional[str] = None
    confidence: float = 0.0
    sub_questions: List[str] = []

# Chain of Thought response
class CoTResponse(BaseModel):
    query: str
    reasoning_type: ReasoningType
    thought_process: List[CoTStep]
    final_answer: str
    confidence_score: float
    evidence_sources: List[str]

# Model Context Protocol message structure
class MCPMessage(BaseModel):
    sender: str
    receiver: str
    type: str
    trace_id: str
    payload: Dict[str, Any]
    timestamp: Optional[str] = None

# Initialize components
embedding_model = SentenceTransformer(EMBEDDING_MODEL)
chroma_client = chromadb.PersistentClient(path=CHROMA_DB_PATH)
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=CHUNK_SIZE,
    chunk_overlap=CHUNK_OVERLAP
)

# Global message queue for MCP communication
message_queue = []

def create_mcp_message(sender: str, receiver: str, msg_type: str, trace_id: str, payload: Dict[str, Any]) -> MCPMessage:
    """Create a standardized MCP message"""
    return MCPMessage(
        sender=sender,
        receiver=receiver,
        type=msg_type,
        trace_id=trace_id,
        payload=payload,
        timestamp=datetime.now().isoformat()
    )

def send_mcp_message(message: MCPMessage):
    """Send message via MCP (add to queue)"""
    message_queue.append(message)
    logger.info(f"MCP Message sent: {message.sender} -> {message.receiver} ({message.type})")

def get_mcp_messages(receiver: str, msg_type: str = None) -> List[MCPMessage]:
    """Retrieve messages for a specific receiver"""
    messages = [msg for msg in message_queue if msg.receiver == receiver]
    if msg_type:
        messages = [msg for msg in messages if msg.type == msg_type]
    return messages

# Document processing utilities
class DocumentProcessor:
    @staticmethod
    def process_pdf(file_path: str) -> str:
        """Extract text from PDF"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text
        except Exception as e:
            logger.error(f"Error processing PDF: {e}")
            return ""

    @staticmethod
    def process_docx(file_path: str) -> str:
        """Extract text from DOCX"""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error processing DOCX: {e}")
            return ""

    @staticmethod
    def process_pptx(file_path: str) -> str:
        """Extract text from PPTX"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            return text
        except Exception as e:
            logger.error(f"Error processing PPTX: {e}")
            return ""

    @staticmethod
    def process_csv(file_path: str) -> str:
        """Extract text from CSV"""
        try:
            df = pd.read_csv(file_path)
            text = f"CSV Data Summary:\n"
            text += f"Columns: {', '.join(df.columns.tolist())}\n"
            text += f"Rows: {len(df)}\n\n"
            text += "Sample Data:\n"
            text += df.head(10).to_string()
            return text
        except Exception as e:
            logger.error(f"Error processing CSV: {e}")
            return ""

    @staticmethod
    def process_txt(file_path: str) -> str:
        """Extract text from TXT/Markdown"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            logger.error(f"Error processing TXT: {e}")
            return ""

# Chain of Thought Reasoning Agent
class CoTReasoningAgent:
    def __init__(self):
        self.name = "CoTReasoningAgent"
        self.ollama_url = f"{OLLAMA_BASE_URL}/api/generate"
    
    def classify_query_complexity(self, query: str) -> ReasoningType:
        """Classify the complexity and type of reasoning needed for a query"""
        query_lower = query.lower()
        
        # Keywords that indicate different reasoning types
        comparison_keywords = ["compare", "contrast", "difference", "similar", "versus", "vs", "better", "worse"]
        analytical_keywords = ["analyze", "explain", "why", "how", "cause", "effect", "reason", "factor"]
        creative_keywords = ["create", "design", "imagine", "brainstorm", "generate", "propose"]
        complex_keywords = ["multiple", "various", "several", "complex", "detailed", "comprehensive"]
        
        if any(keyword in query_lower for keyword in comparison_keywords):
            return ReasoningType.COMPARATIVE
        elif any(keyword in query_lower for keyword in analytical_keywords):
            return ReasoningType.ANALYTICAL
        elif any(keyword in query_lower for keyword in creative_keywords):
            return ReasoningType.CREATIVE
        elif any(keyword in query_lower for keyword in complex_keywords) or len(query.split()) > 15:
            return ReasoningType.COMPLEX
        else:
            return ReasoningType.SIMPLE
    
    def generate_thought_process(self, query: str, context: List[str], reasoning_type: ReasoningType) -> List[CoTStep]:
        """Generate a structured thought process based on query and context"""
        steps = []
        
        if reasoning_type == ReasoningType.SIMPLE:
            steps = self._generate_simple_steps(query, context)
        elif reasoning_type == ReasoningType.COMPLEX:
            steps = self._generate_complex_steps(query, context)
        elif reasoning_type == ReasoningType.ANALYTICAL:
            steps = self._generate_analytical_steps(query, context)
        elif reasoning_type == ReasoningType.COMPARATIVE:
            steps = self._generate_comparative_steps(query, context)
        elif reasoning_type == ReasoningType.CREATIVE:
            steps = self._generate_creative_steps(query, context)
        
        return steps
    
    def _generate_simple_steps(self, query: str, context: List[str]) -> List[CoTStep]:
        """Generate steps for simple queries"""
        steps = []
        
        # Step 1: Understanding the query
        steps.append(CoTStep(
            step_number=1,
            reasoning=f"Understanding the query: '{query}' - This appears to be a straightforward question requiring direct information retrieval.",
            confidence=0.9
        ))
        
        # Step 2: Analyzing available context
        steps.append(CoTStep(
            step_number=2,
            reasoning=f"Analyzing available context: Found {len(context)} relevant document chunks that may contain the answer.",
            evidence=f"Context sources: {len(context)} documents",
            confidence=0.8 if context else 0.3
        ))
        
        # Step 3: Extracting answer
        steps.append(CoTStep(
            step_number=3,
            reasoning="Extracting the most relevant information from the context to directly answer the query.",
            confidence=0.7 if context else 0.2
        ))
        
        return steps
    
    def _generate_complex_steps(self, query: str, context: List[str]) -> List[CoTStep]:
        """Generate steps for complex queries"""
        steps = []
        
        # Step 1: Breaking down the query
        steps.append(CoTStep(
            step_number=1,
            reasoning=f"Breaking down complex query: '{query}' - This requires multiple pieces of information and synthesis.",
            sub_questions=self._extract_sub_questions(query),
            confidence=0.8
        ))
        
        # Step 2: Gathering comprehensive context
        steps.append(CoTStep(
            step_number=2,
            reasoning=f"Gathering comprehensive context: Analyzing {len(context)} document chunks for multiple aspects of the query.",
            evidence=f"Context coverage: {len(context)} sources",
            confidence=0.7 if len(context) >= 3 else 0.4
        ))
        
        # Step 3: Synthesizing information
        steps.append(CoTStep(
            step_number=3,
            reasoning="Synthesizing information from multiple sources to address all aspects of the complex query.",
            confidence=0.6 if context else 0.2
        ))
        
        # Step 4: Validating completeness
        steps.append(CoTStep(
            step_number=4,
            reasoning="Validating that all parts of the query have been addressed and ensuring logical consistency.",
            confidence=0.7
        ))
        
        return steps
    
    def _generate_analytical_steps(self, query: str, context: List[str]) -> List[CoTStep]:
        """Generate steps for analytical queries"""
        steps = []
        
        # Step 1: Identifying analytical requirements
        steps.append(CoTStep(
            step_number=1,
            reasoning=f"Identifying analytical requirements: '{query}' requires explanation of causes, effects, or relationships.",
            confidence=0.8
        ))
        
        # Step 2: Collecting evidence
        steps.append(CoTStep(
            step_number=2,
            reasoning=f"Collecting evidence from {len(context)} sources to support analytical reasoning.",
            evidence=f"Evidence sources: {len(context)} documents",
            confidence=0.7 if context else 0.3
        ))
        
        # Step 3: Analyzing relationships
        steps.append(CoTStep(
            step_number=3,
            reasoning="Analyzing cause-effect relationships, patterns, and underlying mechanisms from the evidence.",
            confidence=0.6 if context else 0.2
        ))
        
        # Step 4: Drawing conclusions
        steps.append(CoTStep(
            step_number=4,
            reasoning="Drawing logical conclusions based on the analysis and supporting them with evidence.",
            confidence=0.7
        ))
        
        return steps
    
    def _generate_comparative_steps(self, query: str, context: List[str]) -> List[CoTStep]:
        """Generate steps for comparative queries"""
        steps = []
        
        # Step 1: Identifying comparison elements
        steps.append(CoTStep(
            step_number=1,
            reasoning=f"Identifying comparison elements in: '{query}' - Determining what aspects need to be compared.",
            confidence=0.8
        ))
        
        # Step 2: Gathering information for each element
        steps.append(CoTStep(
            step_number=2,
            reasoning=f"Gathering information for each comparison element from {len(context)} sources.",
            evidence=f"Comparison data sources: {len(context)} documents",
            confidence=0.7 if len(context) >= 2 else 0.4
        ))
        
        # Step 3: Establishing comparison criteria
        steps.append(CoTStep(
            step_number=3,
            reasoning="Establishing clear criteria for comparison based on the available information.",
            confidence=0.6
        ))
        
        # Step 4: Making comparisons
        steps.append(CoTStep(
            step_number=4,
            reasoning="Making systematic comparisons across the established criteria.",
            confidence=0.7 if context else 0.3
        ))
        
        return steps
    
    def _generate_creative_steps(self, query: str, context: List[str]) -> List[CoTStep]:
        """Generate steps for creative queries"""
        steps = []
        
        # Step 1: Understanding creative requirements
        steps.append(CoTStep(
            step_number=1,
            reasoning=f"Understanding creative requirements: '{query}' requires generating new ideas or solutions.",
            confidence=0.7
        ))
        
        # Step 2: Gathering inspiration from context
        steps.append(CoTStep(
            step_number=2,
            reasoning=f"Gathering inspiration and relevant information from {len(context)} sources.",
            evidence=f"Inspiration sources: {len(context)} documents",
            confidence=0.6 if context else 0.4
        ))
        
        # Step 3: Brainstorming and ideation
        steps.append(CoTStep(
            step_number=3,
            reasoning="Brainstorming ideas and solutions based on the available information and creative requirements.",
            confidence=0.5
        ))
        
        # Step 4: Refining and presenting ideas
        steps.append(CoTStep(
            step_number=4,
            reasoning="Refining the generated ideas and presenting them in a structured, actionable format.",
            confidence=0.6
        ))
        
        return steps
    
    def _extract_sub_questions(self, query: str) -> List[str]:
        """Extract potential sub-questions from a complex query"""
        # Simple heuristic to identify sub-questions
        sub_questions = []
        
        # Split by common conjunctions
        parts = query.replace(" and ", " | ").replace(" or ", " | ").replace(", ", " | ").split(" | ")
        
        for part in parts:
            part = part.strip()
            if len(part) > 10:  # Filter out very short fragments
                if not part.endswith("?"):
                    part += "?"
                sub_questions.append(part)
        
        return sub_questions[:5]  # Limit to 5 sub-questions
    
    def calculate_overall_confidence(self, steps: List[CoTStep], has_context: bool) -> float:
        """Calculate overall confidence score based on reasoning steps"""
        if not steps:
            return 0.0
        
        # Base confidence from steps
        step_confidence = sum(step.confidence for step in steps) / len(steps)
        
        # Adjust based on context availability
        context_factor = 0.8 if has_context else 0.4
        
        # Combine factors
        overall_confidence = (step_confidence * 0.7) + (context_factor * 0.3)
        
        return min(overall_confidence, 1.0)

# Agent 1: IngestionAgent
class IngestionAgent:
    def __init__(self):
        self.name = "IngestionAgent"
        self.collection_name = "documents"
        self.processor = DocumentProcessor()
        
        # Initialize ChromaDB collection
        try:
            self.collection = chroma_client.get_collection(self.collection_name)
        except:
            self.collection = chroma_client.create_collection(
                name=self.collection_name,
                metadata={"hnsw:space": "cosine"}
            )

    def process_document(self, file_path: str, filename: str) -> List[LangChainDocument]:
        """Process document and return chunks"""
        file_ext = Path(filename).suffix.lower()
        
        if file_ext == '.pdf':
            text = self.processor.process_pdf(file_path)
        elif file_ext == '.docx':
            text = self.processor.process_docx(file_path)
        elif file_ext == '.pptx':
            text = self.processor.process_pptx(file_path)
        elif file_ext == '.csv':
            text = self.processor.process_csv(file_path)
        elif file_ext in ['.txt', '.md']:
            text = self.processor.process_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
        
        if not text.strip():
            raise ValueError("No text extracted from document")
        
        # Create document and split into chunks
        doc = LangChainDocument(page_content=text, metadata={"source": filename})
        chunks = text_splitter.split_documents([doc])
        
        return chunks

    def embed_and_store(self, chunks: List[LangChainDocument], trace_id: str):
        """Embed chunks and store in ChromaDB"""
        texts = [chunk.page_content for chunk in chunks]
        embeddings = embedding_model.encode(texts).tolist()
        
        # Generate IDs for chunks
        ids = [f"{trace_id}_chunk_{i}" for i in range(len(chunks))]
        
        # Prepare metadata
        metadatas = []
        for i, chunk in enumerate(chunks):
            metadata = chunk.metadata.copy()
            metadata.update({
                "chunk_id": i,
                "trace_id": trace_id,
                "timestamp": datetime.now().isoformat()
            })
            metadatas.append(metadata)
        
        # Store in ChromaDB
        self.collection.add(
            ids=ids,
            embeddings=embeddings,
            documents=texts,
            metadatas=metadatas
        )
        
        logger.info(f"Stored {len(chunks)} chunks in ChromaDB")
        return len(chunks)

# Agent 2: RetrievalAgent
class RetrievalAgent:
    def __init__(self):
        self.name = "RetrievalAgent"
        self.collection_name = "documents"
        
        # Get ChromaDB collection
        try:
            self.collection = chroma_client.get_collection(self.collection_name)
        except:
            raise ValueError("No document collection found. Please ingest documents first.")

    def retrieve_context(self, query: str, top_k: int = TOP_K_RETRIEVAL) -> Tuple[List[str], List[str]]:
        """Retrieve relevant context for a query"""
        # Generate query embedding
        query_embedding = embedding_model.encode([query]).tolist()
        
        # Search in ChromaDB
        results = self.collection.query(
            query_embeddings=query_embedding,
            n_results=top_k,
            include=["documents", "metadatas"]
        )
        
        # Extract relevant documents
        documents = results["documents"][0] if results["documents"] else []
        metadatas = results["metadatas"][0] if results["metadatas"] else []
        
        # Format context with source information
        contexts = []
        sources = []
        for doc, metadata in zip(documents, metadatas):
            source = metadata.get("source", "unknown")
            context = f"Source: {source}\nContent: {doc}"
            contexts.append(context)
            sources.append(source)
        
        return contexts, sources

# Agent 3: Enhanced LLMResponseAgent with Chain of Thought
class LLMResponseAgent:
    def __init__(self):
        self.name = "LLMResponseAgent"
        self.ollama_url = f"{OLLAMA_BASE_URL}/api/generate"
        self.cot_agent = CoTReasoningAgent()

    def generate_response_with_cot(self, query: str, context: List[str], sources: List[str]) -> CoTResponse:
        """Generate response using Chain of Thought reasoning"""
        
        # Step 1: Classify query complexity
        reasoning_type = self.cot_agent.classify_query_complexity(query)
        
        # Step 2: Generate thought process
        thought_process = self.cot_agent.generate_thought_process(query, context, reasoning_type)
        
        # Step 3: Generate the actual response using LLM
        final_answer = self._generate_llm_response(query, context, thought_process)
        
        # Step 4: Calculate confidence
        confidence_score = self.cot_agent.calculate_overall_confidence(thought_process, bool(context))
        
        return CoTResponse(
            query=query,
            reasoning_type=reasoning_type,
            thought_process=thought_process,
            final_answer=final_answer,
            confidence_score=confidence_score,
            evidence_sources=list(set(sources))
        )

    def _generate_llm_response(self, query: str, context: List[str], thought_process: List[CoTStep]) -> str:
        """Generate LLM response with Chain of Thought context"""
        
        # Format context
        formatted_context = "\n\n".join(context)
        
        # Format thought process
        formatted_thoughts = "\n".join([
            f"Step {step.step_number}: {step.reasoning}" for step in thought_process
        ])
        
        # Create enhanced prompt with Chain of Thought
        prompt = f"""You are a helpful AI assistant that uses step-by-step reasoning to answer questions accurately.

My reasoning process for this query:
{formatted_thoughts}

Context from documents:
{formatted_context}

Question: {query}

Based on my reasoning process and the available context, here is my answer:"""

        # Prepare request payload
        payload = {
            "model": "llama3.2:3b",
            "prompt": prompt,
            "stream": False,
            "options": {
                "temperature": 0.7,
                "max_tokens": 800
            }
        }
        
        try:
            response = requests.post(self.ollama_url, json=payload, timeout=30)
            response.raise_for_status()
            
            result = response.json()
            return result.get("response", "No response generated")
            
        except requests.RequestException as e:
            logger.error(f"Error calling Ollama: {e}")
            return f"Error generating response: {str(e)}"

# Initialize agents
ingestion_agent = IngestionAgent()
retrieval_agent = RetrievalAgent()
llm_agent = LLMResponseAgent()

# API Endpoints

@app.get("/")
async def root():
    return {"message": "RAG Chatbot with Agentic Architecture & Chain of Thought", "version": "1.0.0"}

@app.post("/ingest")
async def ingest_document(file: UploadFile = File(...)):
    """Ingest a document via IngestionAgent"""
    trace_id = str(uuid.uuid4())
    
    try:
        # Save uploaded file temporarily
        temp_file_path = f"./temp_{file.filename}"
        with open(temp_file_path, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
        
        # Send MCP message to IngestionAgent
        mcp_message = create_mcp_message(
            sender="API",
            receiver="IngestionAgent",
            msg_type="DOCUMENT_UPLOAD",
            trace_id=trace_id,
            payload={
                "filename": file.filename,
                "file_path": temp_file_path,
                "content_type": file.content_type
            }
        )
        send_mcp_message(mcp_message)
        
        # Process document
        chunks = ingestion_agent.process_document(temp_file_path, file.filename)
        num_chunks = ingestion_agent.embed_and_store(chunks, trace_id)
        
        # Clean up temp file
        os.remove(temp_file_path)
        
        # Send completion message
        completion_message = create_mcp_message(
            sender="IngestionAgent",
            receiver="API",
            msg_type="INGESTION_COMPLETE",
            trace_id=trace_id,
            payload={
                "filename": file.filename,
                "chunks_created": num_chunks,
                "status": "success"
            }
        )
        send_mcp_message(completion_message)
        
        return {
            "message": "Document ingested successfully",
            "trace_id": trace_id,
            "filename": file.filename,
            "chunks_created": num_chunks
        }
        
    except Exception as e:
        logger.error(f"Error ingesting document: {e}")
        # Clean up temp file if it exists
        if os.path.exists(temp_file_path):
            os.remove(temp_file_path)
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/query")
async def query_documents(query: str = Form(...)):
    """Query documents with Chain of Thought reasoning"""
    trace_id = str(uuid.uuid4())
    
    try:
        # Send query to RetrievalAgent
        retrieval_message = create_mcp_message(
            sender="API",
            receiver="RetrievalAgent",
            msg_type="QUERY_REQUEST",
            trace_id=trace_id,
            payload={"query": query}
        )
        send_mcp_message(retrieval_message)
        
        # Retrieve context
        context, sources = retrieval_agent.retrieve_context(query)
        
        # Send retrieved context to LLMResponseAgent with CoT
        llm_message = create_mcp_message(
            sender="RetrievalAgent",
            receiver="LLMResponseAgent",
            msg_type="RETRIEVAL_RESULT",
            trace_id=trace_id,
            payload={
                "retrieved_context": context,
                "sources": sources,
                "query": query
            }
        )
        send_mcp_message(llm_message)
        
        # Generate response with Chain of Thought
        cot_response = llm_agent.generate_response_with_cot(query, context, sources)
        
        # Send final response
        response_message = create_mcp_message(
            sender="LLMResponseAgent",
            receiver="API",
            msg_type="COT_RESPONSE_GENERATED",
            trace_id=trace_id,
            payload={
                "query": query,
                "response": cot_response.dict(),
                "context_count": len(context)
            }
        )
        send_mcp_message(response_message)
        
        return {
            "query": query,
            "reasoning_type": cot_response.reasoning_type.value,
            "thought_process": [step.dict() for step in cot_response.thought_process],
            "final_answer": cot_response.final_answer,
            "confidence_score": cot_response.confidence_score,
            "evidence_sources": cot_response.evidence_sources,
            "trace_id": trace_id,
            "context_count": len(context)
        }
        
    except Exception as e:
        logger.error(f"Error processing query: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/query-simple")
async def query_documents_simple(query: str = Form(...)):
    """Simple query endpoint without Chain of Thought (for comparison)"""
    trace_id = str(uuid.uuid4())
    
    try:
        # Retrieve context
        context, sources = retrieval_agent.retrieve_context(query)
        
        # Generate simple response
        formatted_context = "\n\n".join(context)
        prompt = f"""You are a helpful assistant. Use the context below to answer the question.

Context:
{formatted_context}

Question: {query}

Answer:"""

        payload = {
            "model": "llama3.2:3b",
            "prompt": prompt,
            "stream": False,
            "options": {
                "temperature": 0.7,
                "max_tokens": 500
            }
        }
        
        response = requests.post(f"{OLLAMA_BASE_URL}/api/generate", json=payload, timeout=30)
        response.raise_for_status()
        result = response.json()
        simple_response = result.get("response", "No response generated")
        
        return {
            "query": query,
            "response": simple_response,
            "trace_id": trace_id,
            "context_count": len(context),
            "sources": list(set(sources))
        }
        
    except Exception as e:
        logger.error(f"Error processing simple query: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/health")
async def health_check():
    """Health check endpoint"""
    try:
        # Check ChromaDB
        collections = chroma_client.list_collections()
        chroma_status = "healthy"
    except Exception as e:
        chroma_status = f"error: {str(e)}"
    
    try:
        # Check Ollama
        ollama_response = requests.get(f"{OLLAMA_BASE_URL}/api/tags", timeout=5)
        ollama_status = "healthy" if ollama_response.status_code == 200 else "unhealthy"
    except Exception as e:
        ollama_status = f"error: {str(e)}"
    
    return {
        "status": "healthy",
        "components": {
            "chromadb": chroma_status,
            "ollama": ollama_status,
            "embedding_model": EMBEDDING_MODEL,
            "chain_of_thought": "enabled"
        }
    }

@app.delete("/reset")
async def reset_database():
    """Reset the vector database and reinitialize agents"""
    global ingestion_agent, retrieval_agent, llm_agent, message_queue

    try:
        # Delete and recreate the collection
        chroma_client.delete_collection(name="documents")
        chroma_client.create_collection(
            name="documents",
            metadata={"hnsw:space": "cosine"}
        )

        # Clear message queue
        message_queue.clear()

        # Reinitialize agents
        ingestion_agent = IngestionAgent()
        retrieval_agent = RetrievalAgent()
        llm_agent = LLMResponseAgent()

        logger.info("ChromaDB and agents reset successfully")
        return {"message": "Database and agents reset successfully"}
        
    except Exception as e:
        logger.error(f"Error resetting database: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/reasoning-stats")
async def get_reasoning_stats():
    """Get statistics about reasoning patterns from message queue"""
    try:
        # Filter CoT response messages
        cot_messages = [msg for msg in message_queue if msg.type == "COT_RESPONSE_GENERATED"]
        
        if not cot_messages:
            return {"message": "No Chain of Thought queries processed yet"}
        
        # Analyze reasoning patterns
        reasoning_types = {}
        confidence_scores = []
        
        for msg in cot_messages:
            response_data = msg.payload.get("response", {})
            if isinstance(response_data, dict):
                reasoning_type = response_data.get("reasoning_type", "unknown")
                confidence = response_data.get("confidence_score", 0.0)
                
                reasoning_types[reasoning_type] = reasoning_types.get(reasoning_type, 0) + 1
                confidence_scores.append(confidence)
        
        avg_confidence = sum(confidence_scores) / len(confidence_scores) if confidence_scores else 0.0
        
        return {
            "total_cot_queries": len(cot_messages),
            "reasoning_type_distribution": reasoning_types,
            "average_confidence": round(avg_confidence, 3),
            "confidence_range": {
                "min": min(confidence_scores) if confidence_scores else 0,
                "max": max(confidence_scores) if confidence_scores else 0
            }
        }
        
    except Exception as e:
        logger.error(f"Error getting reasoning stats: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/mcp-messages")
async def get_mcp_messages(receiver: str = None, msg_type: str = None):
    """Get MCP messages for debugging and monitoring"""
    try:
        filtered_messages = message_queue
        
        if receiver:
            filtered_messages = [msg for msg in filtered_messages if msg.receiver == receiver]
        
        if msg_type:
            filtered_messages = [msg for msg in filtered_messages if msg.type == msg_type]
        
        # Convert to dict for JSON serialization
        messages_dict = []
        for msg in filtered_messages[-50:]:  # Return last 50 messages
            messages_dict.append({
                "sender": msg.sender,
                "receiver": msg.receiver,
                "type": msg.type,
                "trace_id": msg.trace_id,
                "timestamp": msg.timestamp,
                "payload_keys": list(msg.payload.keys()) if msg.payload else []
            })
        
        return {
            "total_messages": len(message_queue),
            "filtered_messages": len(filtered_messages),
            "messages": messages_dict
        }
        
    except Exception as e:
        logger.error(f"Error getting MCP messages: {e}")
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    import uvicorn
    
    # Create necessary directories
    os.makedirs(CHROMA_DB_PATH, exist_ok=True)
    
    print("🚀 Starting RAG Chatbot with Agentic Architecture & Chain of Thought")
    print("📋 Requirements:")
    print("  - Ollama running on http://localhost:11434")
    print("  - Model 'llama3.2:3b' available in Ollama")
    print("🔧 To install Ollama model: ollama pull llama3.2:3b")
    print("🌐 API will be available at: http://localhost:8000")
    print("📚 API docs at: http://localhost:8000/docs")
    print("🧠 Chain of Thought reasoning enabled")
    print("🔍 New endpoints:")
    print("  - /query (with CoT reasoning)")
    print("  - /query-simple (without CoT)")
    print("  - /reasoning-stats (CoT analytics)")
    print("  - /mcp-messages (debugging)")
    
    uvicorn.run(app, host="0.0.0.0", port=8000)